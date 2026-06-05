from pathlib import Path

import docx
import openpyxl
import pdfplumber
from pptx import Presentation

from app.schemas.documents import ExtractedFragment
from app.services.ocr import OcrService


class ExtractionService:
    def __init__(self, ocr: OcrService | None = None) -> None:
        self.ocr = ocr or OcrService()

    def extract(self, path: str, source_file: str) -> list[ExtractedFragment]:
        suffix = Path(path).suffix.lower()

        if suffix == ".docx":
            return self._extract_docx(path, source_file)
        if suffix == ".xlsx":
            return self._extract_xlsx(path, source_file)
        if suffix == ".pdf":
            return self._extract_pdf(path, source_file)
        if suffix == ".pptx":
            return self._extract_pptx(path, source_file)
        if suffix in {".png", ".jpg", ".jpeg", ".bmp", ".webp"}:
            return self.ocr.extract_image(path, source_file)
        return []

    def _extract_docx(self, path: str, source_file: str) -> list[ExtractedFragment]:
        document = docx.Document(path)
        fragments: list[ExtractedFragment] = []

        for index, paragraph in enumerate(document.paragraphs, start=1):
            text = paragraph.text.strip()
            if text:
                fragments.append(
                    ExtractedFragment(
                        text=text,
                        source_file=source_file,
                        source_location=f"paragraph {index}",
                        confidence=1.0,
                    )
                )

        for table_index, table in enumerate(document.tables, start=1):
            for row_index, row in enumerate(table.rows, start=1):
                values = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if values:
                    fragments.append(
                        ExtractedFragment(
                            text=" | ".join(values),
                            source_file=source_file,
                            source_location=f"table {table_index} row {row_index}",
                            confidence=1.0,
                            kind="table",
                        )
                    )

        return fragments

    def _extract_xlsx(self, path: str, source_file: str) -> list[ExtractedFragment]:
        workbook = openpyxl.load_workbook(path, data_only=True)
        fragments: list[ExtractedFragment] = []

        try:
            for worksheet in workbook.worksheets:
                for row in worksheet.iter_rows(values_only=True):
                    values = [str(value) for value in row if value is not None]
                    if values:
                        fragments.append(
                            ExtractedFragment(
                                text=" | ".join(values),
                                source_file=source_file,
                                source_location=f"sheet {worksheet.title}",
                                confidence=1.0,
                                kind="table",
                            )
                        )
        finally:
            workbook.close()

        return fragments

    def _extract_pdf(self, path: str, source_file: str) -> list[ExtractedFragment]:
        fragments: list[ExtractedFragment] = []

        with pdfplumber.open(path) as pdf:
            for index, page in enumerate(pdf.pages, start=1):
                text = (page.extract_text() or "").strip()
                if text:
                    fragments.append(
                        ExtractedFragment(
                            text=text,
                            source_file=source_file,
                            source_location=f"page {index}",
                            confidence=0.95,
                        )
                    )

        return fragments

    def _extract_pptx(self, path: str, source_file: str) -> list[ExtractedFragment]:
        presentation = Presentation(path)
        fragments: list[ExtractedFragment] = []

        for slide_index, slide in enumerate(presentation.slides, start=1):
            for shape_index, shape in enumerate(slide.shapes, start=1):
                text = getattr(shape, "text", "").strip()
                if text:
                    fragments.append(
                        ExtractedFragment(
                            text=text,
                            source_file=source_file,
                            source_location=f"slide {slide_index} shape {shape_index}",
                            confidence=1.0,
                        )
                    )

        return fragments
